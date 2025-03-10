import streamlit as st
import openai
import pandas as pd
import json
from docx import Document
from pptx import Presentation
import PyPDF2
from io import BytesIO
from fpdf import FPDF

# Set your OpenAI API key
openai.api_key = "your_openai_api_key"

def extract_text_from_file(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()
    text = ""
    
    if file_type == "txt":
        text = uploaded_file.read().decode("utf-8")
    elif file_type == "pdf":
        reader = PyPDF2.PdfReader(uploaded_file)
        text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif file_type == "docx":
        doc = Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_type == "pptx":
        presentation = Presentation(uploaded_file)
        text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    
    return text

def chat_with_document(user_input, document_text):
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are an AI assistant that answers questions based on the provided document."},
            {"role": "user", "content": f"Document Text: {document_text}\n\nUser Question: {user_input}"}
        ]
    )
    return response['choices'][0]['message']['content']

def generate_flashcards(text):
    prompt = f"""
    Extract key concepts from the following text and convert them into flashcards.
    Format: JSON with 'question' and 'answer' fields.
    
    Text:
    {text}
    """
    
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "system", "content": "You are a helpful assistant that creates educational flashcards."},
                  {"role": "user", "content": prompt}]
    )
    
    flashcards = json.loads(response['choices'][0]['message']['content'])
    return flashcards

def create_pdf_flashcards(flashcards):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", style='B', size=16)
    pdf.cell(200, 10, "Flashcards", ln=True, align='C')
    pdf.ln(10)
    
    pdf.set_font("Arial", size=12)
    for i, card in enumerate(flashcards):
        pdf.set_font("Arial", style='B', size=12)
        pdf.cell(0, 10, f"Flashcard {i+1}", ln=True)
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 8, f"Q: {card['question']}")
        pdf.multi_cell(0, 8, f"A: {card['answer']}")
        pdf.ln(5)
    
    buffer = BytesIO()
    pdf.output(buffer)
    buffer.seek(0)
    return buffer

st.title("Smart Document Assistant")
st.write("Upload any document (PDF, Word, PPT, TXT), chat with it, and generate flashcards.")

uploaded_file = st.file_uploader("Upload a document", type=["txt", "pdf", "docx", "pptx"]) 
document_text = ""

if uploaded_file:
    document_text = extract_text_from_file(uploaded_file)
    st.subheader("Document Preview")
    st.text_area("Extracted Text", document_text[:1000], height=200)
    
    st.subheader("Chat with Your Document")
    user_query = st.text_input("Ask a question about the document")
    if user_query:
        answer = chat_with_document(user_query, document_text)
        st.write("**Answer:**", answer)
    
    if st.button("Generate Flashcards"):
        flashcards = generate_flashcards(document_text)
        if flashcards:
            st.subheader("Preview Flashcards")
            for i, card in enumerate(flashcards):
                with st.expander(f"Flashcard {i+1}"):
                    st.write(f"**Q:** {card['question']}")
                    st.write(f"**A:** {card['answer']}")
            
            # Convert flashcards to PDF for download
            pdf_flashcards = create_pdf_flashcards(flashcards)
            st.download_button(
                label="Download Flashcards as PDF",
                data=pdf_flashcards,
                file_name="flashcards.pdf",
                mime="application/pdf"
            )
